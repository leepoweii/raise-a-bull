"""Unit tests for document parsers."""
import io
import pytest
from raisebull.parsers.document import parse_docx, parse_xlsx, parse_pptx, parse_pdf


class TestParseDocx:
    def test_simple_docx(self):
        from docx import Document
        doc = Document()
        doc.add_paragraph("Hello World")
        doc.add_paragraph("Second paragraph")
        buf = io.BytesIO()
        doc.save(buf)
        result = parse_docx(buf.getvalue(), "test.docx")
        assert "Hello World" in result
        assert "Second paragraph" in result

    def test_empty_docx(self):
        from docx import Document
        doc = Document()
        buf = io.BytesIO()
        doc.save(buf)
        result = parse_docx(buf.getvalue(), "test.docx")
        assert result == "(空檔案)"

    def test_docx_with_table(self):
        from docx import Document
        doc = Document()
        doc.add_paragraph("Before table")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "H1"
        table.cell(0, 1).text = "H2"
        table.cell(1, 0).text = "A"
        table.cell(1, 1).text = "B"
        buf = io.BytesIO()
        doc.save(buf)
        result = parse_docx(buf.getvalue(), "test.docx")
        assert "Before table" in result
        assert "H1" in result
        assert "|" in result


class TestParseXlsx:
    def test_simple_xlsx(self):
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["Name", "Value"])
        ws.append(["Alice", 10])
        buf = io.BytesIO()
        wb.save(buf)
        result = parse_xlsx(buf.getvalue(), "test.xlsx")
        assert "Sheet: Data" in result
        assert "Alice" in result

    def test_empty_sheet(self):
        import openpyxl
        wb = openpyxl.Workbook()
        buf = io.BytesIO()
        wb.save(buf)
        result = parse_xlsx(buf.getvalue(), "test.xlsx")
        assert "空工作表" in result

    def test_large_xlsx_truncation(self):
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["id", "value"])
        for i in range(200):
            ws.append([i, i * 10])
        buf = io.BytesIO()
        wb.save(buf)
        result = parse_xlsx(buf.getvalue(), "test.xlsx")
        assert "showing first" in result


class TestParsePptx:
    def test_simple_pptx(self):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title Slide"
        slide.placeholders[1].text = "Body text"
        buf = io.BytesIO()
        prs.save(buf)
        result = parse_pptx(buf.getvalue(), "test.pptx")
        assert "Slide 1" in result
        assert "Title Slide" in result

    def test_empty_pptx(self):
        from pptx import Presentation
        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        result = parse_pptx(buf.getvalue(), "test.pptx")
        assert result == "(空檔案)"


class TestParsePdf:
    @pytest.mark.asyncio
    async def test_simple_pdf(self):
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Hello PDF")
        pdf_bytes = doc.tobytes()
        doc.close()
        result = await parse_pdf(pdf_bytes, "test.pdf")
        assert "Hello PDF" in result

    @pytest.mark.asyncio
    async def test_empty_pdf(self):
        # Minimal valid PDF with 0 pages (hand-crafted bytes; fitz.tobytes() refuses
        # to serialize a 0-page document in newer PyMuPDF versions).
        empty_pdf = (
            b"%PDF-1.4\n"
            b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
            b"2 0 obj\n<< /Type /Pages /Kids [] /Count 0 >>\nendobj\n"
            b"xref\n0 3\n"
            b"0000000000 65535 f \n"
            b"0000000009 00000 n \n"
            b"0000000058 00000 n \n"
            b"trailer\n<< /Size 3 /Root 1 0 R >>\n"
            b"startxref\n110\n%%EOF\n"
        )
        result = await parse_pdf(empty_pdf, "test.pdf")
        assert result == "(空檔案)"
