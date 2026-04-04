"""Document file parsers: PDF, DOCX, XLSX, PPTX."""

import io
from typing import Callable, Awaitable, Optional

MAX_PDF_PAGES = 20
MAX_XLSX_ROWS = 100


async def parse_pdf(
    file_bytes: bytes,
    filename: str,
    vision_fn: Optional[Callable[[bytes], Awaitable[str]]] = None,
) -> str:
    """Parse PDF: extract text per page. Falls back to vision_fn for scanned pages.

    Args:
        file_bytes: Raw PDF bytes.
        filename: Original filename (for error context).
        vision_fn: Async callable that accepts image bytes and returns text.
                   Called when a page has < 50 chars of extracted text.

    Returns:
        Pages joined by '\\n\\n---\\n\\n'. Caps at MAX_PDF_PAGES pages.
    """
    import fitz  # PyMuPDF

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        total_pages = len(doc)

        if total_pages == 0:
            return "(空檔案)"

        pages_to_read = min(total_pages, MAX_PDF_PAGES)
        page_texts = []

        for i in range(pages_to_read):
            page = doc[i]
            text = page.get_text().strip()

            if len(text) < 50 and vision_fn is not None:
                # Scanned page — render to image and call vision_fn
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("png")
                text = await vision_fn(img_bytes)

            page_texts.append(text if text else "(no text)")

        result = "\n\n---\n\n".join(page_texts)
        if total_pages > MAX_PDF_PAGES:
            result += f"\n\n*(showing {MAX_PDF_PAGES} of {total_pages} pages)*"

        return result
    finally:
        doc.close()


def parse_docx(file_bytes: bytes, filename: str) -> str:
    """Parse DOCX: extract paragraphs and tables as markdown.

    Returns:
        Paragraphs separated by newlines; tables rendered as markdown tables.
    """
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(io.BytesIO(file_bytes))
    parts = []

    # Iterate body elements in document order to preserve paragraph/table order
    body = doc.element.body
    for child in body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "p":
            # Paragraph
            text = "".join(run.text for run in child.findall(f".//{qn('w:r')}/{qn('w:t')}"))
            if text.strip():
                parts.append(text)

        elif tag == "tbl":
            # Table — build markdown
            rows = child.findall(f".//{qn('w:tr')}")
            if not rows:
                continue

            md_rows = []
            for row in rows:
                cells = row.findall(f".//{qn('w:tc')}")
                cell_texts = []
                for cell in cells:
                    cell_text = "".join(
                        t.text or ""
                        for t in cell.findall(f".//{qn('w:t')}")
                    ).replace("|", "\\|")
                    cell_texts.append(cell_text)
                md_rows.append(cell_texts)

            if md_rows:
                header = md_rows[0]
                sep = ["---"] * len(header)
                lines = [
                    "| " + " | ".join(header) + " |",
                    "| " + " | ".join(sep) + " |",
                ]
                for row in md_rows[1:]:
                    lines.append("| " + " | ".join(row) + " |")
                parts.append("\n".join(lines))

    return "\n\n".join(parts) if parts else "(空檔案)"


def parse_xlsx(file_bytes: bytes, filename: str) -> str:
    """Parse XLSX: each sheet as a markdown table. Caps at MAX_XLSX_ROWS per sheet.

    Returns:
        Sheets separated by '\\n\\n---\\n\\n', each with a '## Sheet: <name>' header.
    """
    import openpyxl

    wb = openpyxl.load_workbook(
        io.BytesIO(file_bytes), read_only=True, data_only=True
    )
    try:
        sheet_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))

            if not rows:
                sheet_parts.append(f"## Sheet: {sheet_name}\n\n(空工作表)")
                continue

            total_rows = len(rows)
            truncated = total_rows > MAX_XLSX_ROWS
            rows_to_render = rows[:MAX_XLSX_ROWS]

            header = rows_to_render[0]
            data_rows = rows_to_render[1:]

            def cell_str(v):
                return str(v).replace("|", "\\|") if v is not None else ""

            sep = ["---"] * len(header)
            lines = [
                "| " + " | ".join(cell_str(c) for c in header) + " |",
                "| " + " | ".join(sep) + " |",
            ]
            for row in data_rows:
                lines.append("| " + " | ".join(cell_str(c) for c in row) + " |")

            content = "\n".join(lines)
            if truncated:
                content += f"\n\n*(showing first {MAX_XLSX_ROWS} of {total_rows} rows)*"

            sheet_parts.append(f"## Sheet: {sheet_name}\n\n{content}")

        return "\n\n---\n\n".join(sheet_parts) if sheet_parts else "(空檔案)"
    finally:
        wb.close()


def parse_pptx(file_bytes: bytes, filename: str) -> str:
    """Parse PPTX: extract text per slide as '## Slide N'.

    Returns:
        Slides joined by '\\n\\n'.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slide_parts = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)

        content = "\n".join(texts) if texts else "(no text)"
        slide_parts.append(f"## Slide {i}\n\n{content}")

    return "\n\n".join(slide_parts) if slide_parts else "(空檔案)"
